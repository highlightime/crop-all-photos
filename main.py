from PIL import Image
import os
from PyPDF2 import PdfMerger, PdfReader
from pptx import Presentation
from pptx.util import Cm
import argparse

# Implement the Crop function
def crop(args):
    # Perform the Crop function with the specified crop box
    crop_box = [int(args.crop_box_l), int(args.crop_box_t), int(args.crop_box_r), int(args.crop_box_b)]
    crop_images(args.input_dir, args.output_dir, crop_box)

# Implement the PDF function
def pdf(args):
    # Perform the PDF conversion function
    image_to_pdf(args.input_dir, args.output_dir)

# Implement the Merge function
def merge(args):
    # Merge PDF files function
    merge_pdfs(args.input_dir, args.output_file)

# Implement the PPTX function
def pptx(args):
    # Convert images to PPTX function
    images_to_pptx(args.input_dir, args.output_file)

# Function to crop images
def crop_images(input_dir, output_dir, crop_box):
    # Create the output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.mkdir(output_dir)

    # Get all image files in the input directory
    image_files = [f for f in os.listdir(input_dir) if f.endswith((".jpg", ".jpeg", ".png"))]

    # Crop and save all images
    for image_file in image_files:
        input_path = os.path.join(input_dir, image_file)
        img = Image.open(input_path)
        cropped_img = img.crop(crop_box)
        output_path = os.path.join(output_dir, image_file)
        cropped_img.save(output_path)

    print("Image cropping is completed.")

# Function to convert images to PDF
def image_to_pdf(input_dir, output_dir):
    # Create the output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.mkdir(output_dir)

    # Get all image files in the input directory
    image_files = [f for f in os.listdir(input_dir) if f.endswith((".jpg", ".jpeg", ".png"))]

    # Convert images to PDF
    for image_file in image_files:
        input_path = os.path.join(input_dir, image_file)
        output_pdf = os.path.join(output_dir, os.path.splitext(image_file)[0] + ".pdf")
        img = Image.open(input_path)
        img = img.convert('RGB')
        img.save(output_pdf, "PDF")

    print("Each image is converted to PDF.")

# Function to merge PDF files
def merge_pdfs(input_dir, output_pdf):
    # Get all PDF files in the input directory
    pdf_files = [f for f in os.listdir(input_dir) if f.endswith(".pdf")]

    # Merge PDF files
    merger = PdfMerger()
    for pdf_file in pdf_files:
        input_path = os.path.join(input_dir, pdf_file)
        merger.append(PdfReader(open(input_path, "rb")))

    # Save the merged PDF file
    merger.write(output_pdf)
    merger.close()

    print(f"PDF file {output_pdf} is created.")

# Function to convert images to PPTX
def images_to_pptx(image_folder, pptx_file):
    # Create a PPTX presentation
    pptx = Presentation()

    # Get image files in the image folder and sort them by name
    image_files = [f for f in os.listdir(image_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
    image_files.sort()

    for image_file in image_files:
        slide = pptx.slides.add_slide(pptx.slide_layouts[6])  # 6 is the image slide layout.
        left = top = 0
        width = Cm(33)
        height = Cm(19)

        image_path = os.path.join(image_folder, image_file)
        pic = slide.shapes.add_picture(image_path, left, top, width, height)

    # Save the PPTX file
    print(f"PPT file {pptx_file} is created.")
    pptx.save(pptx_file)

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description="Programs that can crop many images at once or make them into PDFs and PPTs")
    subparsers = parser.add_subparsers(title="subcommand", dest="subcommand")

    # Add the Crop subcommand
    crop_parser = subparsers.add_parser("crop", help="Function to crop images")
    crop_parser.add_argument("input_dir", help="Input Folder Name")
    crop_parser.add_argument("output_dir", help="Output Folder Name")
    crop_parser.add_argument("crop_box_l", help="Left point of crop area")
    crop_parser.add_argument("crop_box_t", help="Top point of crop area")
    crop_parser.add_argument("crop_box_r", help="Right point of crop area")
    crop_parser.add_argument("crop_box_b", help="Bottom point of crop area")

    # Add the PDF subcommand
    pdf_parser = subparsers.add_parser("pdf", help="Function to convert images to PDF")
    pdf_parser.add_argument("input_dir", help="Input Folder Name")
    pdf_parser.add_argument("output_dir", help="Output Folder Name")

    # Add the Merge subcommand
    merge_parser = subparsers.add_parser("merge", help="Function to merge PDF files")
    merge_parser.add_argument("input_dir", help="Input Folder Name")
    merge_parser.add_argument("output_file", help="Output File Name")

    # Add the PPTX subcommand
    pptx_parser = subparsers.add_parser("pptx", help="Function to convert images to PPTX")
    pptx_parser.add_argument("input_dir", help="Input Folder Name")
    pptx_parser.add_argument("output_file", help="Output File Name")

    args = parser.parse_args()

    if args.subcommand == "crop":
        crop(args)
    elif args.subcommand == "pdf":
        pdf(args)
    elif args.subcommand == "merge":
        merge(args)
    elif args.subcommand == "pptx":
        pptx(args)
    else:
        print("Invalid Subcommands.")
